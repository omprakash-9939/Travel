# -*- coding: utf-8 -*-
"""
Builds the stakeholder pitch deck for the DataArt Travel
Intent-Based Personalization Engine.
Run: python build_deck.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---------------------------------------------------------------- palette
NAVY      = RGBColor(0x0E, 0x2A, 0x47)   # deep ocean navy  (primary)
NAVY_2    = RGBColor(0x12, 0x37, 0x5C)   # lighter navy
CORAL     = RGBColor(0xFF, 0x6B, 0x5C)   # warm coral       (accent 1)
TEAL      = RGBColor(0x19, 0xC3, 0xB1)   # teal             (accent 2)
AMBER     = RGBColor(0xFF, 0xB3, 0x4D)   # amber
LIGHT     = RGBColor(0xF5, 0xF7, 0xFA)   # light bg
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
INK       = RGBColor(0x1A, 0x23, 0x30)   # near-black text
MUTED     = RGBColor(0x6B, 0x72, 0x80)   # muted grey
CARD      = RGBColor(0xFF, 0xFF, 0xFF)
CARD_DK   = RGBColor(0x18, 0x3A, 0x5F)

FONT_H = "Segoe UI Semibold"
FONT_B = "Segoe UI"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


# ---------------------------------------------------------------- helpers
def slide():
    return prs.slides.add_slide(BLANK)

def bg(s, color):
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = color

def rect(s, x, y, w, h, color, line=None, shape=MSO_SHAPE.RECTANGLE):
    sp = s.shapes.add_shape(shape, x, y, w, h)
    sp.fill.solid()
    sp.fill.fore_color.rgb = color
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(1)
    sp.shadow.inherit = False
    return sp

def no_autofit(tf):
    # prevent text from forcing autosize weirdness
    tf.word_wrap = True

def txt(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
        space_after=6, line_spacing=1.0):
    """runs: list of paragraphs; each paragraph = list of (text, size, color, bold, font)"""
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    no_autofit(tf)
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        p.line_spacing = line_spacing
        for (t, sz, col, bold, fnt) in para:
            r = p.add_run()
            r.text = t
            r.font.size = Pt(sz)
            r.font.color.rgb = col
            r.font.bold = bold
            r.font.name = fnt
    return tb

def P(text, size, color, bold=False, font=FONT_B):
    """single-run paragraph helper"""
    return [(text, size, color, bold, font)]

def chip(s, x, y, text, fill, fg=WHITE, w=None, h=Inches(0.42), size=12):
    if w is None:
        w = Inches(0.22 + 0.11 * len(text))
    sp = rect(s, x, y, w, h, fill, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    sp.adjustments[0] = 0.5
    tf = sp.text_frame
    tf.word_wrap = True
    tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = True
    r.font.color.rgb = fg; r.font.name = FONT_H
    return sp

def accent_bar(s, x, y, w=Inches(0.9), color=CORAL, h=Inches(0.09)):
    rect(s, x, y, w, h, color)

def page_header(s, kicker, title, kicker_color=CORAL):
    """standard light content slide header"""
    txt(s, Inches(0.8), Inches(0.55), Inches(11), Inches(0.4),
        [P(kicker.upper(), 13, kicker_color, True, FONT_H)])
    txt(s, Inches(0.8), Inches(0.92), Inches(11.7), Inches(0.9),
        [P(title, 30, NAVY, True, FONT_H)])
    accent_bar(s, Inches(0.82), Inches(1.62), color=kicker_color)

def footer(s, idx, dark=False):
    col = RGBColor(0xAE,0xBA,0xC8) if not dark else RGBColor(0x6F,0x86,0xA3)
    txt(s, Inches(0.8), Inches(7.06), Inches(8), Inches(0.3),
        [P("DataArt Travel  ·  Smart Travel Personalization", 9, col, False, FONT_B)])
    txt(s, Inches(11.6), Inches(7.06), Inches(1.0), Inches(0.3),
        [P(str(idx), 9, col, False, FONT_B)], align=PP_ALIGN.RIGHT)


# ================================================================ 1. TITLE
def s_title(idx):
    s = slide(); bg(s, NAVY)
    # decorative diagonal accent blocks
    rect(s, Inches(0), Inches(0), Inches(0.28), SH, CORAL)
    rect(s, Inches(10.9), Inches(0), Inches(2.43), SH, NAVY_2)
    rect(s, Inches(10.9), Inches(0), Inches(0.12), SH, TEAL)
    # subtle dotted "flight path"
    for i in range(9):
        d = rect(s, Inches(10.95 + i*0.26), Inches(0.9 + i*0.62),
                 Inches(0.07), Inches(0.07), TEAL, shape=MSO_SHAPE.OVAL)
    rect(s, Inches(12.55), Inches(5.95), Inches(0.34), Inches(0.34), CORAL, shape=MSO_SHAPE.ISOSCELES_TRIANGLE)

    chip(s, Inches(0.85), Inches(1.15), "PRODUCT PITCH  ·  STAKEHOLDER REVIEW", NAVY_2, TEAL, w=Inches(4.4), size=12)
    txt(s, Inches(0.8), Inches(1.95), Inches(10), Inches(2.4),
        [P("Knowing What", 52, WHITE, True, FONT_H),
         [("the Traveler ", 52, WHITE, True, FONT_H), ("Wants", 52, CORAL, True, FONT_H)],
         [("— before they ask.", 52, TEAL, True, FONT_H)]],
        line_spacing=1.02)
    txt(s, Inches(0.82), Inches(4.95), Inches(9.6), Inches(1.0),
        [P("A smart engine that learns what each traveler wants — and turns searches into bookings.",
           18, RGBColor(0xC9,0xD6,0xE3), False, FONT_B)], line_spacing=1.15)
    rect(s, Inches(0.85), Inches(5.95), Inches(3.4), Inches(0.045), CORAL)
    txt(s, Inches(0.82), Inches(6.15), Inches(9), Inches(0.8),
        [[("DataArt Travel  ", 14, WHITE, True, FONT_H),
          ("·  MERN Personalization Squad  ·  June 2026", 14, RGBColor(0x9F,0xB1,0xC4), False, FONT_B)]])
    return s


# ================================================================ 2. OPENING STORY
def s_story(idx):
    s = slide(); bg(s, LIGHT)
    rect(s, Inches(0), Inches(0), SW, Inches(0.18), CORAL)
    txt(s, Inches(0.8), Inches(0.6), Inches(11), Inches(0.4),
        [P("MEET MAYA", 13, CORAL, True, FONT_H)])
    txt(s, Inches(0.8), Inches(0.95), Inches(11.7), Inches(0.7),
        [P("She is ready to book. We just can't see her.", 30, NAVY, True, FONT_H)])
    accent_bar(s, Inches(0.82), Inches(1.62))

    steps = [
        ("Mon", "Maya searches Delhi -> Barcelona. Looks at 6 flights.", "She leaves to check prices elsewhere."),
        ("Wed", "She comes back. Searches Barcelona again. Saves a hotel.", "Her interest is clearly growing."),
        ("Today", "She opens a rival app first — they remembered her.", "We showed her the same page as everyone."),
    ]
    x = Inches(0.8); y = Inches(2.05); cw = Inches(3.83); gap = Inches(0.13)
    for i,(day, line, sub) in enumerate(steps):
        cx = Emu(int(x) + i*(int(cw)+int(gap)))
        card = rect(s, cx, y, cw, Inches(3.0), CARD, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        card.adjustments[0] = 0.06
        col = [CORAL, AMBER, NAVY][i]
        chip(s, Emu(int(cx)+int(Inches(0.3))), Emu(int(y)+int(Inches(0.3))), day, col, WHITE, w=Inches(1.2), size=13)
        txt(s, Emu(int(cx)+int(Inches(0.3))), Emu(int(y)+int(Inches(1.0))), Inches(3.25), Inches(1.4),
            [P(line, 15, INK, True, FONT_H)], line_spacing=1.12)
        txt(s, Emu(int(cx)+int(Inches(0.3))), Emu(int(y)+int(Inches(2.25))), Inches(3.25), Inches(0.6),
            [P(sub, 13, MUTED, False, FONT_B)], line_spacing=1.1)

    band = rect(s, Inches(0.8), Inches(5.35), Inches(11.73), Inches(1.15), NAVY, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    band.adjustments[0] = 0.12
    txt(s, Inches(1.2), Inches(5.58), Inches(11.0), Inches(0.8),
        [[("The signs were there the whole time. ", 18, WHITE, True, FONT_H),
          ("Maya isn't a lost cause — she's a missed ", 18, RGBColor(0xC9,0xD6,0xE3), False, FONT_B),
          ("chance", 18, TEAL, True, FONT_H),
          (".", 18, RGBColor(0xC9,0xD6,0xE3), False, FONT_B)]],
        anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.1)
    footer(s, idx)
    return s


# ================================================================ 3. PROBLEM
def s_problem(idx):
    s = slide(); bg(s, WHITE)
    page_header(s, "The Problem", "Ready-to-book travelers leave — and we never learn why")
    txt(s, Inches(0.8), Inches(1.85), Inches(11.7), Inches(0.6),
        [P("People who clearly want to book end up leaving. Right now we have no way to spot them, win them back, or learn from them.",
           15, MUTED, False, FONT_B)], line_spacing=1.15)

    cards = [
        ("We can't see early visitors", "We don't track people before they log in. So we miss most of what they do on the site.", CORAL),
        ("Right message, wrong time", "Our reminders can arrive up to 2 hours late — long after the person has moved on.", AMBER),
        ("We risk losing trust", "\"Prices dropped 7%\" alerts are made up, not real. That can break trust and even cause legal trouble.", TEAL),
        ("We don't measure results", "We can't tell if a reminder led to a booking. Without numbers, we can't prove anything works.", NAVY_2),
    ]
    x0 = Inches(0.8); y0 = Inches(2.7); cw = Inches(5.78); ch = Inches(1.83); gx = Inches(0.17); gy = Inches(0.2)
    for i,(t, d, col) in enumerate(cards):
        r = i // 2; c = i % 2
        cx = Emu(int(x0) + c*(int(cw)+int(gx)))
        cy = Emu(int(y0) + r*(int(ch)+int(gy)))
        card = rect(s, cx, cy, cw, ch, LIGHT, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        card.adjustments[0] = 0.07
        rect(s, cx, Emu(int(cy)+int(Inches(0.22))), Inches(0.09), Inches(1.4), col)
        txt(s, Emu(int(cx)+int(Inches(0.35))), Emu(int(cy)+int(Inches(0.22))), Inches(5.2), Inches(0.5),
            [P(t, 18, NAVY, True, FONT_H)])
        txt(s, Emu(int(cx)+int(Inches(0.35))), Emu(int(cy)+int(Inches(0.78))), Inches(5.15), Inches(1.0),
            [P(d, 12.5, INK, False, FONT_B)], line_spacing=1.12)
    footer(s, idx)
    return s


# ================================================================ 4. COST OF INACTION
def s_cost(idx):
    s = slide(); bg(s, NAVY)
    txt(s, Inches(0.8), Inches(0.6), Inches(11), Inches(0.4),
        [P("WHY IT MATTERS", 13, TEAL, True, FONT_H)])
    txt(s, Inches(0.8), Inches(0.97), Inches(11.7), Inches(0.7),
        [P("Every visitor we ignore is money walking out the door", 28, WHITE, True, FONT_H)])
    accent_bar(s, Inches(0.82), Inches(1.66), color=CORAL)

    stats = [
        ("~98%", "of people who search never book on a normal travel site. That's our 2% starting point.", CORAL),
        ("+8–18%", "more bookings reported when travel sites win back returning visitors.", TEAL),
        ("Rs.1.7–2.8L", "extra money each month at our middle estimate (2,000 searches a month).", AMBER),
    ]
    x0 = Inches(0.8); y = Inches(2.15); cw = Inches(3.83); gap = Inches(0.12)
    for i,(big, sub, col) in enumerate(stats):
        cx = Emu(int(x0) + i*(int(cw)+int(gap)))
        card = rect(s, cx, y, cw, Inches(2.55), CARD_DK, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        card.adjustments[0] = 0.07
        rect(s, Emu(int(cx)+int(Inches(0.32))), Emu(int(y)+int(Inches(0.32))), Inches(0.7), Inches(0.09), col)
        txt(s, Emu(int(cx)+int(Inches(0.3))), Emu(int(y)+int(Inches(0.55))), Inches(3.3), Inches(0.9),
            [P(big, 40, col, True, FONT_H)])
        txt(s, Emu(int(cx)+int(Inches(0.32))), Emu(int(y)+int(Inches(1.5))), Inches(3.2), Inches(0.95),
            [P(sub, 13, RGBColor(0xC9,0xD6,0xE3), False, FONT_B)], line_spacing=1.13)

    txt(s, Inches(0.8), Inches(5.1), Inches(11.7), Inches(1.6),
        [[("The system already exists in our code. ", 16, WHITE, True, FONT_H),
          ("What's missing is the work to fix the data, measure the results, and act at the right moment. That's exactly what this product does.",
           16, RGBColor(0xB8,0xC7,0xD6), False, FONT_B)]], line_spacing=1.2)
    txt(s, Inches(0.8), Inches(6.55), Inches(11.7), Inches(0.4),
        [P("These numbers are industry averages (McKinsey 2023, Salesforce 2024), not our own data yet. We will replace them with our own.",
           9.5, RGBColor(0x7C,0x91,0xA8), False, FONT_B)])
    return s


# ================================================================ 5. SOLUTION
def s_solution(idx):
    s = slide(); bg(s, WHITE)
    page_header(s, "Our Solution", "A simple engine that listens, learns, and acts", kicker_color=TEAL)
    txt(s, Inches(0.8), Inches(1.85), Inches(11.7), Inches(0.6),
        [P("Every action becomes a clue. Clues add up to an interest score. The score tells us the right thing to do — at the right time.",
           15, MUTED, False, FONT_B)], line_spacing=1.15)

    pills = [
        ("1", "LISTEN", "We watch searches, views, saved items, and bookings as they happen.", CORAL),
        ("2", "SCORE", "We turn those actions into a 0–100 score: low, medium, or high.", AMBER),
        ("3", "DECIDE", "We find where they want to go and pick the right next step.", TEAL),
        ("4", "ACT", "We send a reminder, a \"Continue Planning\" card, or a follow-up email.", NAVY_2),
    ]
    x0 = Inches(0.8); y = Inches(2.75); cw = Inches(2.86); gap = Inches(0.13)
    for i,(n, t, d, col) in enumerate(pills):
        cx = Emu(int(x0) + i*(int(cw)+int(gap)))
        card = rect(s, cx, y, cw, Inches(3.0), LIGHT, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        card.adjustments[0] = 0.06
        circ = rect(s, Emu(int(cx)+int(Inches(0.3))), Emu(int(y)+int(Inches(0.32))), Inches(0.7), Inches(0.7), col, shape=MSO_SHAPE.OVAL)
        ctf = circ.text_frame; ctf.paragraphs[0].alignment = PP_ALIGN.CENTER
        rr = ctf.paragraphs[0].add_run(); rr.text = n; rr.font.size = Pt(22); rr.font.bold = True
        rr.font.color.rgb = WHITE; rr.font.name = FONT_H
        txt(s, Emu(int(cx)+int(Inches(0.3))), Emu(int(y)+int(Inches(1.2))), Inches(2.3), Inches(0.5),
            [P(t, 17, NAVY, True, FONT_H)])
        txt(s, Emu(int(cx)+int(Inches(0.3))), Emu(int(y)+int(Inches(1.72))), Inches(2.32), Inches(1.2),
            [P(d, 12.5, INK, False, FONT_B)], line_spacing=1.13)
        if i < 3:
            ar = rect(s, Emu(int(cx)+int(cw)-int(Inches(0.02))), Emu(int(y)+int(Inches(1.25))),
                      Inches(0.18), Inches(0.18), col, shape=MSO_SHAPE.CHEVRON)
    footer(s, idx)
    return s


# ================================================================ 6. HOW IT WORKS
def s_how(idx):
    s = slide(); bg(s, LIGHT)
    page_header(s, "How It Works", "From a click to a booking — how it flows")

    # pipeline boxes
    nodes = [
        ("User action", "search · view · save · book", CORAL),
        ("Activity tracker", "records each action + spots repeat searches", AMBER),
        ("Interest score", "0–100, level, where they want to go", TEAL),
        ("Decision", "score + no spam + test group", NAVY_2),
        ("Action", "reminder · card · email", NAVY),
    ]
    x = Inches(0.8); y = Inches(2.1); bw = Inches(2.18); bh = Inches(1.35); gap = Inches(0.21)
    for i,(t, d, col) in enumerate(nodes):
        cx = Emu(int(x) + i*(int(bw)+int(gap)))
        card = rect(s, cx, y, bw, bh, WHITE, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        card.adjustments[0] = 0.09
        rect(s, cx, y, bw, Inches(0.13), col)
        txt(s, Emu(int(cx)+int(Inches(0.18))), Emu(int(y)+int(Inches(0.28))), Inches(1.9), Inches(0.5),
            [P(t, 14.5, NAVY, True, FONT_H)], line_spacing=1.0)
        txt(s, Emu(int(cx)+int(Inches(0.18))), Emu(int(y)+int(Inches(0.74))), Inches(1.92), Inches(0.6),
            [P(d, 10.5, MUTED, False, FONT_B)], line_spacing=1.08)
        if i < 4:
            rect(s, Emu(int(cx)+int(bw)+int(Inches(0.015))), Emu(int(y)+int(Inches(0.5))),
                 Inches(0.2), Inches(0.34), TEAL, shape=MSO_SHAPE.CHEVRON)

    # two supporting callouts
    callouts = [
        ("Fast when it matters", "Reminders show up while the person is still on the site — not 2 hours later.", TEAL),
        ("Honest by design", "Fake price alerts stay OFF until we have real price data. No made-up urgency.", CORAL),
        ("Built to be measured", "We hold back 20% of users as a test group, so we can prove it really works.", AMBER),
    ]
    cy = Inches(4.0); cw = Inches(3.83); cgap = Inches(0.12); cx0 = Inches(0.8)
    for i,(t, d, col) in enumerate(callouts):
        cx = Emu(int(cx0) + i*(int(cw)+int(cgap)))
        card = rect(s, cx, cy, cw, Inches(1.95), WHITE, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        card.adjustments[0] = 0.08
        rect(s, Emu(int(cx)+int(Inches(0.28))), Emu(int(cy)+int(Inches(0.28))), Inches(0.55), Inches(0.09), col)
        txt(s, Emu(int(cx)+int(Inches(0.3))), Emu(int(cy)+int(Inches(0.46))), Inches(3.25), Inches(0.5),
            [P(t, 15, NAVY, True, FONT_H)])
        txt(s, Emu(int(cx)+int(Inches(0.3))), Emu(int(cy)+int(Inches(0.98))), Inches(3.25), Inches(0.9),
            [P(d, 12, INK, False, FONT_B)], line_spacing=1.12)
    txt(s, Inches(0.8), Inches(6.18), Inches(11.7), Inches(0.4),
        [P("Built with: React · Express · MongoDB — plus Amadeus, SendGrid, Stripe and Razorpay.",
           11, MUTED, False, FONT_B)])
    footer(s, idx)
    return s


# ================================================================ 7. KEY FEATURES
def s_features(idx):
    s = slide(); bg(s, WHITE)
    page_header(s, "What Travelers Feel", "Four features that do the work behind the scenes", kicker_color=CORAL)
    feats = [
        ("Continue Planning card", "When she returns, the home page shows the exact search she left — no typing again.", CORAL),
        ("Timely reminder", "\"Still planning Barcelona?\" reaches interested users right away — and never spams them.", AMBER),
        ("Smart recommendations", "The home page shows where she wants to go now — not a trip from months ago.", TEAL),
        ("Bring them back by email", "When her interest grows enough, a friendly email brings her back — even after she leaves.", NAVY_2),
    ]
    x0 = Inches(0.8); y0 = Inches(2.0); cw = Inches(5.78); ch = Inches(2.18); gx = Inches(0.17); gy = Inches(0.2)
    for i,(t, d, col) in enumerate(feats):
        r = i//2; c = i%2
        cx = Emu(int(x0)+c*(int(cw)+int(gx)))
        cy = Emu(int(y0)+r*(int(ch)+int(gy)))
        card = rect(s, cx, cy, cw, ch, LIGHT, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        card.adjustments[0] = 0.06
        ic = rect(s, Emu(int(cx)+int(Inches(0.35))), Emu(int(cy)+int(Inches(0.35))), Inches(0.62), Inches(0.62), col, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        ic.adjustments[0]=0.3
        itf = ic.text_frame; itf.paragraphs[0].alignment = PP_ALIGN.CENTER
        ir = itf.paragraphs[0].add_run(); ir.text = "★"; ir.font.size=Pt(20); ir.font.color.rgb=WHITE; ir.font.bold=True
        txt(s, Emu(int(cx)+int(Inches(1.2))), Emu(int(cy)+int(Inches(0.4))), Inches(4.4), Inches(0.7),
            [P(t, 18, NAVY, True, FONT_H)], line_spacing=1.0)
        txt(s, Emu(int(cx)+int(Inches(0.4))), Emu(int(cy)+int(Inches(1.25))), Inches(5.05), Inches(0.85),
            [P(d, 13, INK, False, FONT_B)], line_spacing=1.15)
    footer(s, idx)
    return s


# ================================================================ 8. BUSINESS CASE / ROI
def s_roi(idx):
    s = slide(); bg(s, WHITE)
    page_header(s, "The Business Case", "The numbers work once we have real traffic", kicker_color=TEAL)

    headers = ["Case", "Searches / month", "More bookings", "Extra money / month", "Pays back in"]
    rows = [
        ("Low", "500", "+8%", "Rs. 35k – 70k", "3–8 months"),
        ("Middle", "2,000", "+12%", "Rs. 1.75L – 2.8L", "1–2 months"),
        ("High", "5,000", "+18%", "Rs. 4.2L – 7.0L", "< 1 month"),
    ]
    tx = Inches(0.8); ty = Inches(2.0); tw = Inches(11.73)
    col_w = [Inches(2.5), Inches(2.4), Inches(2.0), Inches(3.13), Inches(1.7)]
    rh = Inches(0.42); rh_body = Inches(0.78)

    # header row
    cx = tx
    for j,h in enumerate(headers):
        cell = rect(s, cx, ty, col_w[j], rh, NAVY)
        tf = cell.text_frame; tf.word_wrap=True; tf.vertical_anchor=MSO_ANCHOR.MIDDLE
        tf.margin_left=Pt(8); tf.margin_right=Pt(6)
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT if j==0 else PP_ALIGN.CENTER
        rr=p.add_run(); rr.text=h; rr.font.size=Pt(12.5); rr.font.bold=True; rr.font.color.rgb=WHITE; rr.font.name=FONT_H
        cx = Emu(int(cx)+int(col_w[j]))
    # body rows
    cy = Emu(int(ty)+int(rh))
    for i,row in enumerate(rows):
        cx = tx
        base = LIGHT if i%2==0 else WHITE
        hl = (i==1)
        for j,val in enumerate(row):
            fill = TEAL if hl and j==0 else (RGBColor(0xE8,0xF8,0xF6) if hl else base)
            cell = rect(s, cx, cy, col_w[j], rh_body, fill)
            cell.line.color.rgb = RGBColor(0xE2,0xE8,0xEF); cell.line.width=Pt(0.75)
            tf=cell.text_frame; tf.word_wrap=True; tf.vertical_anchor=MSO_ANCHOR.MIDDLE
            tf.margin_left=Pt(8); tf.margin_right=Pt(6)
            p=tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT if j==0 else PP_ALIGN.CENTER
            rr=p.add_run(); rr.text=val
            rr.font.size=Pt(13.5 if j==0 else 13)
            rr.font.bold = (j==0) or hl
            rr.font.color.rgb = (WHITE if hl and j==0 else NAVY if j==0 else INK)
            rr.font.name = FONT_H if (j==0 or hl) else FONT_B
            cx=Emu(int(cx)+int(col_w[j]))
        cy=Emu(int(cy)+int(rh_body))

    # takeaway band
    band = rect(s, Inches(0.8), Inches(4.95), Inches(11.73), Inches(1.25), LIGHT, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    band.adjustments[0]=0.1
    rect(s, Inches(0.8), Inches(4.95), Inches(0.12), Inches(1.25), CORAL)
    txt(s, Inches(1.15), Inches(5.15), Inches(11.1), Inches(0.95),
        [[("In short:  ", 15, CORAL, True, FONT_H),
          ("Phase 1 fixes cost only ~3 days and remove legal risk — they pay for themselves many times over. The money case gets strong above ~1,000 searches a month, paying back in 1–2 months in the middle case.",
           14, INK, False, FONT_B)]], anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.15)
    txt(s, Inches(0.8), Inches(6.35), Inches(11.7), Inches(0.4),
        [P("Based on Rs.35k average booking, a 2% starting rate, over 12 months. These are estimates, not promises — we will check them in Phase 2.",
           9.5, MUTED, False, FONT_B)])
    footer(s, idx)
    return s


# ================================================================ 9. ROADMAP
def s_roadmap(idx):
    s = slide(); bg(s, NAVY)
    txt(s, Inches(0.8), Inches(0.6), Inches(11), Inches(0.4),
        [P("HOW WE GET THERE", 13, TEAL, True, FONT_H)])
    txt(s, Inches(0.8), Inches(0.97), Inches(11.7), Inches(0.7),
        [P("A simple step-by-step plan — value before scale", 28, WHITE, True, FONT_H)])
    accent_bar(s, Inches(0.82), Inches(1.66), color=CORAL)

    phases = [
        ("PHASE 1", "Remove Risk", "Week 1 · ~3 days", "Fix the counting bugs, reset after a booking, stop fake price alerts, fix return detection.", CORAL),
        ("PHASE 2", "Measure", "Weeks 2–3 · ~6–9 days", "Add a test group and track clicks, so we know what really works.", AMBER),
        ("PHASE 3", "Deliver Value", "Weeks 4–8 · ~18–29 days", "Faster reminders, smarter recommendations, follow-up emails, a better first visit.", TEAL),
        ("PHASE 4", "Grow", "Month 3+ · ~28–40 days", "Track new visitors (with consent), admin tools, and real live prices.", RGBColor(0x7F,0xB0,0xE8)),
    ]
    # timeline line
    rect(s, Inches(1.0), Inches(2.55), Inches(11.3), Inches(0.05), RGBColor(0x2C,0x4B,0x6E))
    x0 = Inches(0.8); y = Inches(2.95); cw = Inches(2.86); gap = Inches(0.13)
    for i,(ph, t, when, d, col) in enumerate(phases):
        cx = Emu(int(x0)+i*(int(cw)+int(gap)))
        # node dot on the line
        rect(s, Emu(int(cx)+int(Inches(1.25))), Inches(2.42), Inches(0.3), Inches(0.3), col, shape=MSO_SHAPE.OVAL)
        card = rect(s, cx, y, cw, Inches(3.15), CARD_DK, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        card.adjustments[0]=0.06
        chip(s, Emu(int(cx)+int(Inches(0.28))), Emu(int(y)+int(Inches(0.28))), ph, col, NAVY if col==AMBER else WHITE, w=Inches(1.35), size=11)
        txt(s, Emu(int(cx)+int(Inches(0.3))), Emu(int(y)+int(Inches(0.85))), Inches(2.35), Inches(0.5),
            [P(t, 19, WHITE, True, FONT_H)])
        txt(s, Emu(int(cx)+int(Inches(0.3))), Emu(int(y)+int(Inches(1.32))), Inches(2.35), Inches(0.4),
            [P(when, 11.5, col, True, FONT_B)])
        txt(s, Emu(int(cx)+int(Inches(0.3))), Emu(int(y)+int(Inches(1.78))), Inches(2.35), Inches(1.3),
            [P(d, 12, RGBColor(0xC9,0xD6,0xE3), False, FONT_B)], line_spacing=1.13)
    txt(s, Inches(0.8), Inches(6.45), Inches(11.7), Inches(0.4),
        [P("After each step we decide: go or stop. We only spend more once we have proven it works.",
           12, RGBColor(0x9F,0xB1,0xC4), False, FONT_B)])
    return s


# ================================================================ 10. WHY TRUST US
def s_trust(idx):
    s = slide(); bg(s, WHITE)
    page_header(s, "Why Say Yes Now", "We're offering a smart way of working, not just a feature", kicker_color=TEAL)
    points = [
        ("It's already built", "The engine is in our code today. We just need to fix it, measure it, and switch it on — not start over."),
        ("Honest numbers", "Every estimate is marked as an estimate. We promise to replace guesses with our own real data."),
        ("Safety first", "Phase 1 removes the legal risk from fake prices before anything reaches real users."),
        ("Easy to measure", "A built-in test group lets us prove it works — or stop early — instead of guessing."),
    ]
    y = Inches(2.05)
    for i,(t, d) in enumerate(points):
        cy = Emu(int(y)+i*int(Inches(1.12)))
        ic = rect(s, Inches(0.85), cy, Inches(0.62), Inches(0.62), TEAL if i%2==0 else CORAL, shape=MSO_SHAPE.OVAL)
        itf=ic.text_frame; itf.paragraphs[0].alignment=PP_ALIGN.CENTER
        ir=itf.paragraphs[0].add_run(); ir.text="✓"; ir.font.size=Pt(22); ir.font.bold=True; ir.font.color.rgb=WHITE; ir.font.name=FONT_H
        txt(s, Inches(1.75), Emu(int(cy)-int(Inches(0.05))), Inches(10.6), Inches(0.5),
            [P(t, 18, NAVY, True, FONT_H)])
        txt(s, Inches(1.75), Emu(int(cy)+int(Inches(0.42))), Inches(10.6), Inches(0.6),
            [P(d, 13.5, INK, False, FONT_B)], line_spacing=1.1)
    footer(s, idx)
    return s


# ================================================================ 11. TEAM
def s_team(idx):
    s = slide(); bg(s, LIGHT)
    page_header(s, "The Team", "Four people. One shared goal.", kicker_color=CORAL)
    members = [
        ("Product & Research Lead", "Owns the problem and the business case", CORAL, "PL"),
        ("Backend / Data Engineer", "Scoring, background jobs & connections", TEAL, "BE"),
        ("Frontend Engineer", "Home page, recommendations & tracking", AMBER, "FE"),
        ("Testing & Quality Lead", "Testing, measurement & quality", NAVY_2, "QA"),
    ]
    x0=Inches(0.8); y=Inches(2.15); cw=Inches(2.86); gap=Inches(0.13)
    for i,(role, d, col, ini) in enumerate(members):
        cx=Emu(int(x0)+i*(int(cw)+int(gap)))
        card=rect(s, cx, y, cw, Inches(3.6), WHITE, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        card.adjustments[0]=0.05
        rect(s, cx, y, cw, Inches(0.16), col)
        av=rect(s, Emu(int(cx)+int(Inches(0.93))), Emu(int(y)+int(Inches(0.55))), Inches(1.0), Inches(1.0), col, shape=MSO_SHAPE.OVAL)
        atf=av.text_frame; atf.paragraphs[0].alignment=PP_ALIGN.CENTER
        ar=atf.paragraphs[0].add_run(); ar.text=ini; ar.font.size=Pt(26); ar.font.bold=True; ar.font.color.rgb=WHITE; ar.font.name=FONT_H
        txt(s, Emu(int(cx)+int(Inches(0.2))), Emu(int(y)+int(Inches(1.75))), Inches(2.46), Inches(0.4),
            [P("Your Name", 15, NAVY, True, FONT_H)], align=PP_ALIGN.CENTER)
        txt(s, Emu(int(cx)+int(Inches(0.2))), Emu(int(y)+int(Inches(2.2))), Inches(2.46), Inches(0.5),
            [P(role, 13, col, True, FONT_H)], align=PP_ALIGN.CENTER, line_spacing=1.05)
        txt(s, Emu(int(cx)+int(Inches(0.25))), Emu(int(y)+int(Inches(2.85))), Inches(2.36), Inches(0.7),
            [P(d, 11.5, MUTED, False, FONT_B)], align=PP_ALIGN.CENTER, line_spacing=1.12)
    txt(s, Inches(0.8), Inches(6.1), Inches(11.7), Inches(0.4),
        [P("Tip: replace the initials and \"Your Name\" with each teammate's real name and photo before presenting.",
           10.5, MUTED, False, FONT_B)], align=PP_ALIGN.CENTER)
    footer(s, idx)
    return s


# ================================================================ 12. THE ASK / CLOSE
def s_ask(idx):
    s = slide(); bg(s, NAVY)
    rect(s, Inches(0), Inches(0), Inches(0.28), SH, CORAL)
    txt(s, Inches(0.85), Inches(0.7), Inches(11), Inches(0.4),
        [P("THE ASK", 13, TEAL, True, FONT_H)])
    txt(s, Inches(0.82), Inches(1.1), Inches(11.7), Inches(1.0),
        [[("Say yes to ", 34, WHITE, True, FONT_H),
          ("Phase 1 + 2", 34, CORAL, True, FONT_H),
          (" — fix, then measure.", 34, WHITE, True, FONT_H)]])
    accent_bar(s, Inches(0.85), Inches(2.05), color=TEAL)

    asks = [
        ("~9–12 days", "of work to remove risk and set up proper measurement.", CORAL),
        ("1 decision", "30 days of real data, then a clear go / stop choice on Phase 3.", TEAL),
        ("0 downside", "Phase 1 only removes risk — it pays for itself no matter what.", AMBER),
    ]
    x0=Inches(0.85); y=Inches(2.45); cw=Inches(3.8); gap=Inches(0.13)
    for i,(big, sub, col) in enumerate(asks):
        cx=Emu(int(x0)+i*(int(cw)+int(gap)))
        card=rect(s, cx, y, cw, Inches(2.0), CARD_DK, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        card.adjustments[0]=0.07
        txt(s, Emu(int(cx)+int(Inches(0.3))), Emu(int(y)+int(Inches(0.3))), Inches(3.2), Inches(0.6),
            [P(big, 24, col, True, FONT_H)])
        txt(s, Emu(int(cx)+int(Inches(0.3))), Emu(int(y)+int(Inches(0.95))), Inches(3.2), Inches(0.95),
            [P(sub, 13, RGBColor(0xC9,0xD6,0xE3), False, FONT_B)], line_spacing=1.15)

    band=rect(s, Inches(0.85), Inches(4.85), Inches(11.65), Inches(1.55), CARD_DK, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    band.adjustments[0]=0.08
    rect(s, Inches(0.85), Inches(4.85), Inches(0.12), Inches(1.55), TEAL)
    txt(s, Inches(1.2), Inches(5.12), Inches(11.1), Inches(1.1),
        [[("Our vision:  ", 17, TEAL, True, FONT_H),
          ("a DataArt Travel that spots what a traveler wants the moment it shows — so the next Maya books with ",
           16, WHITE, False, FONT_B),
          ("us", 16, CORAL, True, FONT_H),
          (", not someone else. Say yes to Phase 1+2, and we'll bring you the proof.",
           16, WHITE, False, FONT_B)]], anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.18)
    return s


# ================================================================ 13. THANK YOU
def s_thanks(idx):
    s = slide(); bg(s, NAVY)
    rect(s, Inches(10.9), Inches(0), Inches(2.43), SH, NAVY_2)
    rect(s, Inches(10.9), Inches(0), Inches(0.12), SH, CORAL)
    for i in range(9):
        rect(s, Inches(10.95 + i*0.26), Inches(0.9 + i*0.62), Inches(0.07), Inches(0.07), TEAL, shape=MSO_SHAPE.OVAL)
    rect(s, Inches(12.55), Inches(5.95), Inches(0.34), Inches(0.34), CORAL, shape=MSO_SHAPE.ISOSCELES_TRIANGLE)

    txt(s, Inches(0.85), Inches(2.5), Inches(10), Inches(1.2),
        [[("Thank you.", 60, WHITE, True, FONT_H)]])
    rect(s, Inches(0.9), Inches(3.85), Inches(3.4), Inches(0.05), CORAL)
    txt(s, Inches(0.88), Inches(4.1), Inches(9.6), Inches(0.7),
        [P("Questions and doubts are very welcome.", 18, RGBColor(0xC9,0xD6,0xE3), False, FONT_B)])
    txt(s, Inches(0.88), Inches(5.2), Inches(9), Inches(0.5),
        [[("DataArt Travel  ", 14, WHITE, True, FONT_H),
          ("·  MERN Personalization Squad", 14, RGBColor(0x9F,0xB1,0xC4), False, FONT_B)]])
    return s


# ---------------------------------------------------------------- build
builders = [s_title, s_story, s_problem, s_cost, s_solution, s_how,
            s_features, s_roi, s_roadmap, s_trust, s_team, s_ask, s_thanks]
for i, b in enumerate(builders, start=1):
    b(i)

out = "DataArt-Travel-Personalization-Pitch.pptx"
prs.save(out)
print("Saved", out, "with", len(prs.slides._sldIdLst), "slides")
